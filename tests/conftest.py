# tests/conftest.py
import html as _html_module
from pathlib import Path
import pytest
from docx import Document as Docx
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF

@pytest.fixture
def make_docx(tmp_path):
    def _make(name: str, paragraphs: list[str]) -> Path:
        d = Docx()
        for p in paragraphs:
            d.add_paragraph(p)
        out = tmp_path / name
        d.save(out)
        return out
    return _make

@pytest.fixture
def make_pptx(tmp_path):
    def _make(name: str, slide_texts: list[str]) -> Path:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for t in slide_texts:
            slide = prs.slides.add_slide(blank)
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            box.text_frame.text = t
        out = tmp_path / name
        prs.save(out)
        return out
    return _make

@pytest.fixture
def make_pdf(tmp_path):
    def _make(name: str, pages: list[str]) -> Path:
        # Use fitz.Story + DocumentWriter so Vietnamese Unicode survives the
        # roundtrip. Page.insert_text() only supports the PDF base-14 Latin
        # character set and silently corrupts Vietnamese glyphs.
        out = tmp_path / name
        mb = fitz.paper_rect("a4")
        where = mb + (36, 36, -36, -36)
        writer = fitz.DocumentWriter(str(out))
        for body in pages:
            safe = _html_module.escape(body).replace("\n", "<br>")
            story = fitz.Story(html=f"<p>{safe}</p>")
            dev = writer.begin_page(mb)
            story.place(where)
            story.draw(dev, None)
            writer.end_page()
        writer.close()
        return out
    return _make

@pytest.fixture
def make_pdf_blocks(tmp_path):
    """Build a multi-page A4 PDF with text placed in fixed regions per page.

    Each page is a dict with optional "header" (top margin), "body" (middle),
    and "footer" (bottom margin). Lets tests put real running headers/footers
    at true margin y-positions on A4 (595x842).
    """
    regions = {            # (y0, y1) in PDF points; A4 height 842
        "header": (20, 70),    # top band  (< 0.12*842 = 101)
        "body": (130, 720),    # safely outside both margins
        "footer": (785, 830),  # bottom band (> 0.88*842 = 741)
    }

    def _make(name: str, pages: list[dict]) -> Path:
        out = tmp_path / name
        mb = fitz.paper_rect("a4")
        writer = fitz.DocumentWriter(str(out))
        for pg in pages:
            dev = writer.begin_page(mb)
            for key, (y0, y1) in regions.items():
                text = pg.get(key)
                if not text:
                    continue
                safe = _html_module.escape(text).replace("\n", "<br>")
                story = fitz.Story(html=f"<p>{safe}</p>")
                story.place(fitz.Rect(40, y0, mb.x1 - 40, y1))
                story.draw(dev, None)
            writer.end_page()
        writer.close()
        return out
    return _make
