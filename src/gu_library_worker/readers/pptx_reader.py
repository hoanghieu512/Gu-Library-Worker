# src/gu_library_worker/readers/pptx_reader.py
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from gu_library_worker.schema import Unit
from .base import Extraction

def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    parts.append(line)
    return "\n".join(parts)

def read_pptx(path: Path) -> Extraction:
    prs = Presentation(str(path))
    units: list[Unit] = []
    for i, slide in enumerate(prs.slides, start=1):
        text = _slide_text(slide) or f"Slide {i}"  # never empty
        units.append(Unit(type="slide", label=f"Slide {i}",
                          path=[], text=text, page=i))
    return Extraction(kind="slide", units=units)
